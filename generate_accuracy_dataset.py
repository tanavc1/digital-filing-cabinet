import os
import random
from reportlab.pdfgen import canvas
from pptx import Presentation
from docx import Document

# Create directory
DATA_DIR = "accuracy_benchmark_data"
os.makedirs(DATA_DIR, exist_ok=True)

def generate_pdf():
    """Generates complex_rental_agreement.pdf (5 pages)"""
    print("Generating complex_rental_agreement.pdf...")
    c = canvas.Canvas(os.path.join(DATA_DIR, "complex_rental_agreement.pdf"))
    
    # Page 1-2: Filler
    for i in range(1, 3):
        c.drawString(100, 750, f"Page {i}: General Terms")
        c.drawString(100, 700, "1.1 The Tenant agrees to lease the premises...")
        c.showPage()
    
    # Page 3: The Facts + Distractor
    c.drawString(100, 750, "Page 3: Deposits and Fees")
    c.drawString(100, 700, "3.1 Security Deposit")
    c.drawString(100, 680, "The Tenant shall pay a refundable Security Deposit of $2,000 upon signing.")
    
    c.drawString(100, 640, "3.2 Pet Policy")
    c.drawString(100, 620, "Animals are permitted subject to approval.")
    # The Needle
    c.drawString(100, 600, "The Tenant must pay a pet deposit of $500 per animal, non-refundable.")
    
    c.showPage()
    
    # Page 4-5: More Filler
    for i in range(4, 6):
        c.drawString(100, 750, f"Page {i}: Signatures")
        c.drawString(100, 700, "Signed: __________________________")
        c.showPage()
        
    c.save()

def generate_docx():
    """Generates financial_report_Q3.docx"""
    print("Generating financial_report_Q3.docx...")
    doc = Document()
    doc.add_heading('Financial Report Q3 2024', 0)
    
    doc.add_heading('1. Executive Summary', level=1)
    doc.add_paragraph('The quarter showed strong top-line growth despite market headwinds.')
    
    doc.add_heading('2. Revenue Performance', level=1)
    # Key Fact 1
    doc.add_paragraph('Q3 Revenue increased by 14.5% year-over-year to $12.4M.')
    
    doc.add_heading('3. Profitability Analysis', level=1)
    # Synthesis Need (Fact 2 + Reason)
    doc.add_paragraph('However, net profit margin dropped to 8% due to increased R&D spend on the new AI initiative.')
    
    doc.add_paragraph('We expect margins to stabilize in Q4.')
    
    doc.save(os.path.join(DATA_DIR, "financial_report_Q3.docx"))

def generate_pptx():
    """Generates product_roadmap_2025.pptx"""
    print("Generating product_roadmap_2025.pptx...")
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Product Roadmap 2025"
    
    # Slide 2: Phase 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Phase 1: Launch"
    slide.shapes.placeholders[1].text = "Expected Q1 2025\nStatus: On Track"
    
    # Slide 3: Phase 2 (Key Fact)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Phase 2: Expansion"
    # Delay
    slide.shapes.placeholders[1].text = "Phase 2 launch is delayed to Q4 2025 due to chip shortage.\nCritical update."
    
    # Slide 4: Specifications (Needle in footnote-ish area)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Unit Economics"
    body = slide.shapes.placeholders[1]
    body.text = "Target margin: 40%.\nProjected cost per unit: $150."
    
    prs.save(os.path.join(DATA_DIR, "product_roadmap_2025.pptx"))

def generate_txt():
    """Generates needle_history.txt (Large file)"""
    print("Generating needle_history.txt...")
    with open(os.path.join(DATA_DIR, "needle_history.txt"), "w") as f:
        f.write("HISTORY OF THE ZOG EMPIRE\n\n")
        
        # Filler text
        filler = "The empire flourished for many centuries. " * 50 + "\n"
        for _ in range(20):
            f.write(filler)
            
        # The Needle
        f.write("\nIn the year 402, the Great Theft occurred. The lost scepter was hidden in the Caves of Zog.\n")
        
        # More filler
        for _ in range(20):
            f.write(filler)

if __name__ == "__main__":
    generate_pdf()
    generate_docx()
    generate_pptx()
    generate_txt()
    print(f"Dataset generated in {DATA_DIR}")
