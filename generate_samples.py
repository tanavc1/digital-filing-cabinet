
import os
from pptx import Presentation
from pptx.util import Inches

def create_sample_docs():
    output_dir = "sample_docs"
    os.makedirs(output_dir, exist_ok=True)

    # 1. Create a Text File (Topic: Space Exploration - General)
    text_content = """
    Space Exploration: The Next Frontier

    Humanity has always looked to the stars. From the early days of the Apollo missions 
    to the current era of commercial spaceflight, our desire to explore has never waned.

    Key Companies in 2024:
    - SpaceX: Leading the charge with Starship.
    - Blue Origin: Focusing on orbital reefs and lunar landers.
    - NASA: Preparing for Artemis missions to return to the Moon.

    Risks:
    The cost of failure is high, but the potential rewards—multi-planetary life, 
    resource acquisition, and scientific discovery—are immense.
    """
    with open(os.path.join(output_dir, "space_overview.txt"), "w") as f:
        f.write(text_content.strip())
    print(f"Created {output_dir}/space_overview.txt")

    # 2. Create a PPTX File (Topic: Mars Colonization Plan)
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Mars Colonization Strategy"
    subtitle.text = "Q3 2025 Strategic Overview"

    # Content Slide 1
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Mission Goals"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Establish permanent base by 2035"
    p = tf.add_paragraph()
    p.text = "Utilize in-situ resource utilization (ISRU)"
    p = tf.add_paragraph()
    p.text = "Develop food production systems"

    # Content Slide 2 (Financials/Risks)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Budget & Risks"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Estimated Cost: $500 Billion over 10 years"
    p = tf.add_paragraph()
    p.text = "Primary Risk: Radiation exposure during transit"
    
    pptx_path = os.path.join(output_dir, "mars_plan.pptx")
    prs.save(pptx_path)
    print(f"Created {pptx_path}")

if __name__ == "__main__":
    create_sample_docs()
